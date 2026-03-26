from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SNOW_BLUE = RGBColor(0x29, 0xB5, 0xE8)
SNOW_DARK = RGBColor(0x11, 0x27, 0x4E)
SNOW_NAVY = RGBColor(0x05, 0x0A, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0xA8, 0x72)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
GCP_GREEN = RGBColor(0x34, 0xA8, 0x53)
AWS_ORANGE = RGBColor(0xFF, 0x99, 0x00)
AZURE_BLUE = RGBColor(0x00, 0x89, 0xD6)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color, text="", font_size=14, font_color=WHITE, bold=False, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape_type, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
    return shp


def add_text_box(slide, left, top, w, h, text, font_size=18, font_color=SNOW_DARK, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_slide(slide, items, left, top, w, h, font_size=16, font_color=SNOW_DARK):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_arrow(slide, left, top, w, h, color=SNOW_BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SNOW_NAVY)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "Cross-Cloud Declarative Data Sharing",
    font_size=42, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
    "Secure, Programmatic Data Distribution Across GCP, AWS & Azure",
    font_size=24, font_color=SNOW_BLUE, bold=False, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
    "BlackLine  |  Snowflake Solution Engineering",
    font_size=20, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.6),
    "March 2026  |  100% SQL-Based  |  Zero UI Dependencies",
    font_size=16, font_color=GRAY_TEXT, bold=False, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2: Agenda
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Agenda", font_size=36, font_color=SNOW_DARK, bold=True)

add_shape(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

items = [
    "1.  The Challenge: Why Cross-Cloud Data Sharing?",
    "2.  Architecture Overview: Three Clouds, One Workflow",
    "3.  What We Built: End-to-End Walkthrough",
    "4.  RBAC & Security Model",
    "5.  Declarative Sharing Deep Dive (TYPE=DATA)",
    "6.  Cross-Cloud Auto-Fulfillment",
    "7.  Live Demo: SQL Execution Flow",
    "8.  Consumer Experience",
    "9.  Future Enhancements & Roadmap",
    "10. Q&A"
]
add_bullet_slide(slide, items, Inches(1.2), Inches(1.8), Inches(10), Inches(5), font_size=20)

# ============================================================================
# SLIDE 3: The Challenge
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "The Challenge", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
    "Enterprises operate across multiple clouds, but data stays siloed",
    font_size=20, font_color=GRAY_TEXT)

challenges = [
    ("Data Silos", "Financial data on GCP, operations on AWS,\nanalytics on Azure - no unified view"),
    ("ETL Complexity", "Traditional sharing requires pipelines,\nS3 buckets, storage accounts, and cron jobs"),
    ("Security Gaps", "Copying data across clouds means\nmore attack surface and compliance risk"),
    ("Stale Data", "Batch ETL means consumers work\nwith hours or days-old data"),
]

for i, (title, desc) in enumerate(challenges):
    x = Inches(0.8 + (i * 3.1))
    add_shape(slide, x, Inches(2.8), Inches(2.8), Inches(3.5), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.6),
        title, font_size=20, font_color=SNOW_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.7), Inches(2.4), Inches(2.2),
        desc, font_size=15, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 4: Architecture Overview (the key diagram slide)
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
    "Architecture: Three Clouds, One Workflow", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

provider_box = add_shape(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(5.5), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.6), Inches(0.5),
    "GCP Provider (us-central1)", font_size=18, font_color=GCP_GREEN, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.8), Inches(2.3), Inches(3.4), Inches(0.7), GCP_GREEN,
    "SHARED_FINANCIAL_DATA", font_size=13, font_color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(3.0), Inches(3.4), Inches(1.0),
    "JOURNAL_ENTRIES | ACCOUNT_BALANCES\nRECONCILIATIONS | CLOSE_TASKS\nINTERCOMPANY_TRANSACTIONS",
    font_size=11, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.8), Inches(4.1), Inches(3.4), Inches(0.7), GCP_GREEN,
    "SHARED_OPERATIONS_DATA", font_size=13, font_color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(3.4), Inches(1.0),
    "CUSTOMERS | ORDERS | PRODUCTS\nSUPPORT_TICKETS | USAGE_METRICS",
    font_size=11, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(1.2), Inches(6.0), Inches(1.2), Inches(0.6), SNOW_DARK,
    "FINANCIAL_\nDATA_PKG", font_size=9, font_color=WHITE, bold=True)
add_shape(slide, Inches(2.7), Inches(6.0), Inches(1.2), Inches(0.6), SNOW_DARK,
    "OPERATIONS_\nDATA_PKG", font_size=9, font_color=WHITE, bold=True)

middle_x = Inches(5.0)
add_shape(slide, middle_x, Inches(1.5), Inches(3.3), Inches(5.5), RGBColor(0xE3, 0xF2, 0xFD))
add_text_box(slide, middle_x + Inches(0.1), Inches(1.6), Inches(3.1), Inches(0.5),
    "Snowflake Global Services", font_size=16, font_color=SNOW_BLUE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, middle_x + Inches(0.3), Inches(2.3), Inches(2.7), Inches(0.7), SNOW_BLUE,
    "Private Listings", font_size=14, font_color=WHITE, bold=True)
add_text_box(slide, middle_x + Inches(0.1), Inches(3.1), Inches(3.1), Inches(0.5),
    "FINANCIAL_DATA_LISTING\nOPERATIONS_DATA_LISTING",
    font_size=12, font_color=SNOW_DARK, align=PP_ALIGN.CENTER)

add_shape(slide, middle_x + Inches(0.3), Inches(4.0), Inches(2.7), Inches(0.7), SNOW_BLUE,
    "Cross-Cloud Auto-Fulfillment", font_size=13, font_color=WHITE, bold=True)
add_text_box(slide, middle_x + Inches(0.1), Inches(4.8), Inches(3.1), Inches(0.8),
    "SUB_DATABASE_WITH_\nREFERENCE_USAGE",
    font_size=12, font_color=SNOW_DARK, align=PP_ALIGN.CENTER)

add_shape(slide, middle_x + Inches(0.3), Inches(5.8), Inches(2.7), Inches(0.7), SNOW_BLUE,
    "manifest.yml + RBAC", font_size=14, font_color=WHITE, bold=True)

aws_x = Inches(8.8)
add_shape(slide, aws_x, Inches(1.5), Inches(4), Inches(2.3), RGBColor(0xFF, 0xF3, 0xE0))
add_text_box(slide, aws_x + Inches(0.2), Inches(1.6), Inches(3.6), Inches(0.5),
    "AWS Consumer (us-west-2)", font_size=18, font_color=AWS_ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, aws_x + Inches(0.3), Inches(2.3), Inches(3.4), Inches(0.7), AWS_ORANGE,
    "FINANCIAL_DATA_APP", font_size=14, font_color=WHITE, bold=True)
add_text_box(slide, aws_x + Inches(0.2), Inches(3.1), Inches(3.6), Inches(0.5),
    "5 tables - live, zero-copy read access",
    font_size=12, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_shape(slide, aws_x, Inches(4.3), Inches(4), Inches(2.3), RGBColor(0xE1, 0xF5, 0xFE))
add_text_box(slide, aws_x + Inches(0.2), Inches(4.4), Inches(3.6), Inches(0.5),
    "Azure Consumer (centralus)", font_size=18, font_color=AZURE_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide, aws_x + Inches(0.3), Inches(5.1), Inches(3.4), Inches(0.7), AZURE_BLUE,
    "OPERATIONS_DATA_APP", font_size=14, font_color=WHITE, bold=True)
add_text_box(slide, aws_x + Inches(0.2), Inches(5.9), Inches(3.6), Inches(0.5),
    "5 tables - live, zero-copy read access",
    font_size=12, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 5: What We Built
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "What We Built: Complete SQL Workflow", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

steps = [
    ("Step 0", "Org Setup", "Enable ORGADMIN\nGrant MANAGE LISTING\nAUTO FULFILLMENT", "00_org_enable_\nauto_fulfillment.sql"),
    ("Step 1", "RBAC Setup", "Create DATA_SHARING_ADMIN\nLeast-privilege grants\nCreate SHARING_WH", "00_provider_\nrbac_setup.sql"),
    ("Step 2", "Source Data", "2 databases, 10 tables\nFinancial + Operations\nRealistic sample data", "01_create_\ndatabases.sql"),
    ("Step 3", "App Packages", "TYPE=DATA packages\nmanifest.yml upload\nRelease LIVE version", "02/03_create_\napp_package.sql"),
    ("Step 4", "Listings", "Private listings w/\nauto_fulfillment config\nPublish to consumers", "04_create_\nlistings.sql"),
    ("Step 5", "Consumer Install", "CREATE APPLICATION\nFROM LISTING\nQuery shared data", "01_install_\napp.sql"),
]

for i, (step, title, desc, file) in enumerate(steps):
    x = Inches(0.4 + (i * 2.1))
    color = GREEN if i < 5 else ORANGE

    add_shape(slide, x, Inches(1.8), Inches(1.9), Inches(0.6), color,
        step, font_size=14, font_color=WHITE, bold=True)

    add_shape(slide, x, Inches(2.5), Inches(1.9), Inches(0.6), SNOW_DARK,
        title, font_size=13, font_color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.1), Inches(3.3), Inches(1.7), Inches(1.5),
        desc, font_size=12, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

    add_shape(slide, x, Inches(5.0), Inches(1.9), Inches(0.8), LIGHT_GRAY,
        file, font_size=10, font_color=SNOW_DARK, bold=False)

    if i < 5:
        add_arrow(slide, x + Inches(1.95), Inches(2.6), Inches(0.15), Inches(0.3), SNOW_BLUE)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
    "Green = Completed & Verified    |    Orange = Pending Auto-Fulfillment Replication",
    font_size=14, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 6: RBAC & Security Model
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "RBAC & Security: Least-Privilege Model", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(5), Inches(0.5),
    "Provider Side (GCP)", font_size=20, font_color=SNOW_DARK, bold=True)

roles_provider = [
    ("ORGADMIN", "MANAGE LISTING AUTO FULFILLMENT\nSYSTEM$ENABLE_GLOBAL_DATA_SHARING", RGBColor(0xB7, 0x1C, 0x1C)),
    ("ACCOUNTADMIN", "Listing creation with auto_fulfillment\nGrant roles to service users", RGBColor(0xE6, 0x51, 0x00)),
    ("DATA_SHARING_ADMIN", "CREATE APPLICATION PACKAGE\nCREATE LISTING, CREATE DATABASE\nDay-to-day sharing operations", SNOW_BLUE),
]

for i, (role, desc, color) in enumerate(roles_provider):
    y = Inches(2.5 + i * 1.5)
    add_shape(slide, Inches(1.0), y, Inches(2.2), Inches(0.7), color,
        role, font_size=13, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(3.4), y, Inches(2.8), Inches(1.0),
        desc, font_size=12, font_color=GRAY_TEXT)

add_shape(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.2), LIGHT_GRAY)
add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.5),
    "Consumer Side (AWS / Azure)", font_size=20, font_color=SNOW_DARK, bold=True)

roles_consumer = [
    ("ACCOUNTADMIN", "Install applications from listings\nCREATE APPLICATION FROM LISTING", RGBColor(0xE6, 0x51, 0x00)),
    ("app_admin", "Full access to all shared objects\nDefined in manifest.yml", SNOW_BLUE),
    ("app_user", "Read-only access to shared data\nDefined in manifest.yml", GREEN),
]

for i, (role, desc, color) in enumerate(roles_consumer):
    y = Inches(2.5 + i * 1.5)
    add_shape(slide, Inches(7.2), y, Inches(2.2), Inches(0.7), color,
        role, font_size=13, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(9.6), y, Inches(2.8), Inches(1.0),
        desc, font_size=12, font_color=GRAY_TEXT)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.5),
    "Authentication: RSA Key-Pair (SNOWFLAKE_JWT)  |  Service Users: TYPE=SERVICE  |  Zero passwords in scripts",
    font_size=14, font_color=GRAY_TEXT, bold=False, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 7: Declarative Sharing Deep Dive
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Declarative Sharing: TYPE=DATA Deep Dive", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_shape(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.3), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(0.5),
    "How It Works", font_size=22, font_color=SNOW_DARK, bold=True)

how_items = [
    "1. CREATE APPLICATION PACKAGE ... TYPE = DATA",
    "2. PUT manifest.yml to snow://package/<PKG>/versions/LIVE/",
    "3. ALTER APPLICATION PACKAGE ... RELEASE LIVE VERSION",
    "4. CREATE EXTERNAL LISTING ... APPLICATION PACKAGE ...",
    "5. ALTER LISTING ... PUBLISH",
    "6. Consumer: CREATE APPLICATION ... FROM LISTING '<ID>'"
]
add_bullet_slide(slide, how_items, Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0), font_size=15, font_color=SNOW_DARK)

add_shape(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.3), LIGHT_GRAY)
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
    "Key Properties", font_size=22, font_color=SNOW_DARK, bold=True)

props = [
    "No setup script required (vs Native Apps)",
    "No code execution on consumer side",
    "Data stays in provider's storage (zero-copy)",
    "manifest.yml defines roles & shared objects",
    "LIVE version auto-created (no BUILD needed)",
    "Cross-cloud via auto-fulfillment replication",
    "Consumer gets read-only app with RBAC roles"
]
add_bullet_slide(slide, props, Inches(7.2), Inches(2.4), Inches(5.2), Inches(4.0), font_size=15, font_color=SNOW_DARK)

# ============================================================================
# SLIDE 8: Cross-Cloud Auto-Fulfillment
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Cross-Cloud Auto-Fulfillment", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
    "How Snowflake automatically replicates data across cloud providers and regions",
    font_size=18, font_color=GRAY_TEXT)

prereqs = [
    "ORGADMIN calls SYSTEM$ENABLE_GLOBAL_DATA_SHARING_FOR_ACCOUNT",
    "ORGADMIN grants MANAGE LISTING AUTO FULFILLMENT ON ACCOUNT",
    "Listing YAML includes auto_fulfillment.refresh_type: SUB_DATABASE_WITH_REFERENCE_USAGE",
    "Target accounts specified in listing targets.accounts array",
]

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(0.5),
    "Prerequisites (One-Time Setup)", font_size=20, font_color=SNOW_DARK, bold=True)
add_bullet_slide(slide, prereqs, Inches(1.0), Inches(3.1), Inches(5.3), Inches(2.5), font_size=14, font_color=GRAY_TEXT)

flow = [
    "Provider publishes listing with auto_fulfillment config",
    "Snowflake creates a replica of the app package in target region",
    "Underlying database data replicates to consumer's region",
    "Consumer sees is_ready_for_import = true in SHOW AVAILABLE LISTINGS",
    "Consumer can now CREATE APPLICATION FROM LISTING",
    "Ongoing changes auto-replicate (incremental sync)",
]

add_text_box(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(0.5),
    "Replication Flow", font_size=20, font_color=SNOW_DARK, bold=True)
add_bullet_slide(slide, flow, Inches(7.2), Inches(3.1), Inches(5.3), Inches(3.5), font_size=14, font_color=GRAY_TEXT)

add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), RGBColor(0xFF, 0xF8, 0xE1))
add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(1.0),
    "Timing: Initial cross-cloud replication typically takes 30-90 minutes depending on data volume and region distance. "
    "Subsequent incremental syncs are much faster. Monitor with: SHOW AVAILABLE LISTINGS + check is_ready_for_import flag.",
    font_size=14, font_color=RGBColor(0x8D, 0x6E, 0x00))

# ============================================================================
# SLIDE 9: Data Model
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Data Model: 2 Databases, 10 Tables", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_shape(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.3), RGBColor(0xE8, 0xF5, 0xE9))
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(5.4), Inches(0.5),
    "SHARED_FINANCIAL_DATA.CORE", font_size=20, font_color=GCP_GREEN, bold=True)
add_text_box(slide, Inches(0.7), Inches(2.3), Inches(5.4), Inches(0.4),
    "Shared with AWS Consumer  |  26 + 25 + 16 + 15 + 16 rows", font_size=13, font_color=GRAY_TEXT)

fin_tables = [
    ("JOURNAL_ENTRIES", "GL journal entries with debit/credit, entity, currency, status"),
    ("ACCOUNT_BALANCES", "Period-end balances by account, entity, and currency"),
    ("RECONCILIATIONS", "Account reconciliation records with GL vs sub-ledger diffs"),
    ("CLOSE_TASKS", "Period-close task tracking with assignments and priorities"),
    ("INTERCOMPANY_TRANSACTIONS", "IC eliminations with from/to entity and matching status"),
]

for i, (tbl, desc) in enumerate(fin_tables):
    y = Inches(2.9 + i * 0.75)
    add_shape(slide, Inches(0.7), y, Inches(2.5), Inches(0.55), GCP_GREEN,
        tbl, font_size=11, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(3.3), y, Inches(2.8), Inches(0.55),
        desc, font_size=10, font_color=GRAY_TEXT)

add_shape(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(5.3), RGBColor(0xE1, 0xF5, 0xFE))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5.4), Inches(0.5),
    "SHARED_OPERATIONS_DATA.CORE", font_size=20, font_color=AZURE_BLUE, bold=True)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5.4), Inches(0.4),
    "Shared with Azure Consumer  |  15 + 17 + 10 + 15 + 30 rows", font_size=13, font_color=GRAY_TEXT)

ops_tables = [
    ("CUSTOMERS", "Customer master: industry, region, tier, ARR, account manager"),
    ("ORDERS", "Sales orders with product, quantity, amount, ship status"),
    ("PRODUCTS", "Product catalog with pricing, margins, and active status"),
    ("SUPPORT_TICKETS", "Customer support: priority, SLA, category, resolution"),
    ("USAGE_METRICS", "Monthly telemetry: active users, API calls, storage, features"),
]

for i, (tbl, desc) in enumerate(ops_tables):
    y = Inches(2.9 + i * 0.75)
    add_shape(slide, Inches(7.2), y, Inches(2.5), Inches(0.55), AZURE_BLUE,
        tbl, font_size=11, font_color=WHITE, bold=True)
    add_text_box(slide, Inches(9.8), y, Inches(2.8), Inches(0.55),
        desc, font_size=10, font_color=GRAY_TEXT)

# ============================================================================
# SLIDE 10: Consumer Experience
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Consumer Experience", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
    "What the consumer sees and does - it's just 2 SQL statements",
    font_size=20, font_color=GRAY_TEXT)

consumer_steps = [
    "-- Step 1: Check if listing is available",
    "SHOW AVAILABLE LISTINGS;",
    "-- Check is_ready_for_import = true",
    "",
    "-- Step 2: Install the application",
    "CREATE APPLICATION FINANCIAL_DATA_APP",
    "  FROM LISTING '<LISTING_GLOBAL_NAME>';",
    "",
    "-- Step 3: Query the shared data",
    "SELECT * FROM FINANCIAL_DATA_APP.CORE.JOURNAL_ENTRIES;",
    "",
    "-- That's it. Zero ETL. Zero pipelines.",
]

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(7), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(consumer_steps):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(15)
    p.font.name = "Courier New"
    if line.startswith("--"):
        p.font.color.rgb = GREEN
    elif line == "":
        p.font.size = Pt(8)
    else:
        p.font.color.rgb = SNOW_DARK

benefits = [
    "No data copying - zero storage cost",
    "Always current - auto-sync from provider",
    "Role-based access via manifest.yml",
    "No pipelines to build or maintain",
    "Works across GCP/AWS/Azure",
]

add_shape(slide, Inches(8.5), Inches(2.3), Inches(4.3), Inches(4.5), LIGHT_GRAY)
add_text_box(slide, Inches(8.7), Inches(2.4), Inches(3.9), Inches(0.5),
    "Consumer Benefits", font_size=20, font_color=SNOW_DARK, bold=True)
add_bullet_slide(slide, benefits, Inches(8.9), Inches(3.0), Inches(3.7), Inches(3.5), font_size=15, font_color=GRAY_TEXT)

# ============================================================================
# SLIDE 11: File Structure
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Project Structure: Fully Reproducible", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

tree = """DataSharing/
+-- setup/                              # One-time org-level setup
|   +-- 00_org_enable_auto_fulfillment.sql   # ORGADMIN grants
|   +-- 01_gcp_service_user.sql              # Service user + RSA key
|   +-- 02_azure_service_user.sql            # Azure service user
|
+-- gcp_provider/                       # Provider-side scripts (run in order)
|   +-- 00_provider_rbac_setup.sql           # DATA_SHARING_ADMIN role
|   +-- 01_create_databases.sql              # 2 DBs, 10 tables, sample data
|   +-- 02_create_app_package_aws.sql        # FINANCIAL_DATA_PKG (TYPE=DATA)
|   +-- 03_create_app_package_azure.sql      # OPERATIONS_DATA_PKG (TYPE=DATA)
|   +-- 04_create_listings.sql               # Private listings + publish
|   +-- manifests/
|       +-- financial_data_pkg/manifest.yml  # Roles & shared objects
|       +-- operations_data_pkg/manifest.yml
|
+-- aws_consumer/                       # AWS consumer install
|   +-- 01_install_app.sql                   # CREATE APPLICATION FROM LISTING
|
+-- azure_consumer/                     # Azure consumer install
    +-- 01_install_app.sql                   # CREATE APPLICATION FROM LISTING"""

txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11), Inches(5.3))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(tree.split("\n")):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.name = "Courier New"
    if "#" in line:
        p.font.color.rgb = GREEN
    else:
        p.font.color.rgb = SNOW_DARK

# ============================================================================
# SLIDE 12: Future Enhancements
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Future Enhancements & Roadmap", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

phases = [
    ("Phase 1: Current", GREEN, [
        "2 databases, 10 tables shared",
        "Cross-cloud GCP -> AWS + Azure",
        "Private listings with auto-fulfillment",
        "RBAC with custom DATA_SHARING_ADMIN role",
        "100% SQL, zero UI dependencies",
    ]),
    ("Phase 2: Near-Term", SNOW_BLUE, [
        "Data masking policies (column-level security)",
        "Row-access policies for tenant isolation",
        "Secure views for computed/aggregated data",
        "Event tables for audit/lineage tracking",
        "Automated testing & CI/CD pipeline",
    ]),
    ("Phase 3: Advanced", SNOW_DARK, [
        "Bi-directional sharing (consumers -> provider)",
        "Snowflake Streams + Tasks for CDC",
        "Cortex AI-powered data quality monitoring",
        "Dynamic Tables for derived datasets",
        "Marketplace public listings (monetization)",
    ]),
]

for i, (phase, color, items) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.7), color,
        phase, font_size=18, font_color=WHITE, bold=True)
    add_bullet_slide(slide, items, x + Inches(0.2), Inches(2.7), Inches(3.4), Inches(4.0),
        font_size=14, font_color=GRAY_TEXT)

# ============================================================================
# SLIDE 13: Key Lessons Learned
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    "Key Lessons & Findings", font_size=36, font_color=SNOW_DARK, bold=True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

lessons = [
    ("Listing Syntax", "CREATE OR REPLACE not valid for listings.\nMust use DROP IF EXISTS + CREATE pattern."),
    ("Auto-Fulfillment", "ORGADMIN must call SYSTEM$ENABLE_GLOBAL_DATA_SHARING\nBEFORE granting MANAGE LISTING AUTO FULFILLMENT."),
    ("LIVE Version", "Never use BUILD or ADD LIVE VERSION.\nUse ALTER APPLICATION PACKAGE ... RELEASE LIVE VERSION."),
    ("Cross-Cloud Timing", "Initial replication takes 30-90 min.\nMonitor with is_ready_for_import flag."),
    ("BCR 2025_03", "Use CREATE LISTING (not CREATE DATA\nEXCHANGE LISTING) on newer accounts."),
    ("Manifest Rules", "File must be named manifest.yml exactly.\nPUT with AUTO_COMPRESS=false required."),
]

for i, (title, desc) in enumerate(lessons):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.7 + row * 2.7)

    add_shape(slide, x, y, Inches(3.8), Inches(2.3), LIGHT_GRAY)
    add_shape(slide, x, y, Inches(3.8), Inches(0.6), SNOW_BLUE,
        title, font_size=16, font_color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.3),
        desc, font_size=13, font_color=GRAY_TEXT, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 14: Q&A / Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SNOW_NAVY)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
    "Questions & Discussion",
    font_size=44, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.5), Inches(3.7), Inches(4.3), Inches(0.04), SNOW_BLUE, shape_type=MSO_SHAPE.RECTANGLE)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
    "Cross-Cloud Declarative Data Sharing  |  BlackLine + Snowflake",
    font_size=22, font_color=SNOW_BLUE, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
    "All scripts available in the DataSharing/ project directory",
    font_size=16, font_color=GRAY_TEXT, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
    "snow sql -c gcp_central  |  snow sql -c myconnection  |  snow sql -c azure_central",
    font_size=14, font_color=GRAY_TEXT, bold=False, align=PP_ALIGN.CENTER)


# ============================================================================
# SAVE
# ============================================================================
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Cross-Cloud_Declarative_Data_Sharing.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
